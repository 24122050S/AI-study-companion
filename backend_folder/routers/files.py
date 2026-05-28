from fastapi import APIRouter, UploadFile, File, Form, HTTPException
import io
from core import supabase, embeddings, text_splitter, call_groq
from langchain_community.vectorstores import SupabaseVectorStore

router = APIRouter()

@router.post("/api/upload")
async def upload_file(file: UploadFile = File(...), user_id: str = Form(...), notebook_id: str = Form(...)): 
    try:
        content = await file.read()
        ext = file.filename.split('.')[-1].lower()
        chunks, metadatas = [], []
        
        if ext == 'pdf':
            import fitz
            pdf_document = fitz.open(stream=content, filetype="pdf")
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                extracted = page.get_text("text") 
                if extracted.strip():
                    page_chunks = text_splitter.split_text(extracted)
                    chunks.extend(page_chunks)
                    metadatas.extend([{"filename": file.filename, "page": page_num + 1} for _ in page_chunks])
                    
        elif ext == 'txt':
            text = content.decode('utf-8', errors='ignore')
            if text.strip():
                page_chunks = text_splitter.split_text(text)
                chunks.extend(page_chunks)
                metadatas.extend([{"filename": file.filename, "page": 1} for _ in page_chunks])
                
        elif ext in ['doc', 'docx']:
            try:
                import docx
            except ImportError:
                return {"status": "error", "message": "Thiếu thư viện python-docx"}
            doc = docx.Document(io.BytesIO(content))
            text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            if text.strip():
                page_chunks = text_splitter.split_text(text)
                chunks.extend(page_chunks)
                metadatas.extend([{"filename": file.filename, "page": 1} for _ in page_chunks])
                
        elif ext in ['ppt', 'pptx']:
            try:
                from pptx import Presentation
            except ImportError:
                return {"status": "error", "message": "Thiếu thư viện python-pptx"}
            ppt = Presentation(io.BytesIO(content))
            for i, slide in enumerate(ppt.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text: slide_text += shape.text + "\n"
                if slide_text.strip():
                    page_chunks = text_splitter.split_text(slide_text)
                    chunks.extend(page_chunks)
                    metadatas.extend([{"filename": file.filename, "page": i + 1} for _ in page_chunks])
                    
        elif ext in ['jpg', 'jpeg', 'png']:
            try:
                import pytesseract
                from PIL import Image
            except ImportError:
                return {"status": "error", "message": "Thiếu thư viện pytesseract Pillow"}
            
            pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
            img = Image.open(io.BytesIO(content))
            text = pytesseract.image_to_string(img, lang='vie+eng')
            if text.strip():
                page_chunks = text_splitter.split_text(text)
                chunks.extend(page_chunks)
                metadatas.extend([{"filename": file.filename, "page": 1} for _ in page_chunks])
                
        else: 
            raise HTTPException(status_code=400, detail=f"Định dạng .{ext} chưa được hỗ trợ!")

        if not chunks: return {"status": "error", "message": "File rỗng hoặc AI không thể đọc được chữ!"}

        for meta in metadatas:
            meta["user_id"] = user_id
            meta["notebook_id"] = int(notebook_id)

        # Đẩy thẳng lên Supabase pgvector
        vector_store = SupabaseVectorStore(
            embedding=embeddings,
            client=supabase,
            table_name="documents",
            query_name="match_documents"
        )
        vector_store.add_texts(chunks, metadatas=metadatas)

        supabase.table("uploaded_files").insert({"user_id": user_id, "filename": file.filename, "notebook_id": int(notebook_id)}).execute()
        return {"status": "success", "message": f"AI đã học xong file {file.filename}!"}
    except Exception as e: 
        print(f"🔥 LỖI UPLOAD CLOUD: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/api/files/{user_id}/{notebook_id}")
async def get_uploaded_files(user_id: str, notebook_id: int):
    res = supabase.table("uploaded_files").select("*").eq("user_id", user_id).eq("notebook_id", notebook_id).order("upload_date", desc=True).execute()
    return [{"id": r['id'], "filename": r['filename'], "date": r['upload_date']} for r in res.data]

@router.delete("/api/files/{file_id}")
async def delete_file_history(file_id: int):
    try:
        res = supabase.table("uploaded_files").select("user_id", "filename", "notebook_id").eq("id", file_id).execute()
        if not res.data:
            raise HTTPException(status_code=404, detail="Không tìm thấy tài liệu này.")
            
        user_id = res.data[0]['user_id']
        filename = res.data[0]['filename']
        notebook_id = res.data[0]['notebook_id'] 
        
        # Xóa trên mây
        supabase.table("documents").delete().eq("metadata->>filename", filename).eq("metadata->>user_id", user_id).eq("metadata->>notebook_id", str(notebook_id)).execute()
        supabase.table("uploaded_files").delete().eq("id", file_id).execute()
        
        return {"status": "success", "message": f"Đã xóa tài liệu {filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi xóa file: {str(e)}")

# 🚀 ĐÃ NÂNG CẤP: Tìm nguồn trích dẫn trực tiếp từ Cloud
@router.get("/api/reference")
async def get_reference_content(user_id: str, filename: str, page: int, notebook_id: str = ""):
    try:
        def normalize_filename(name): return name.replace("_", "").replace(" ", "").replace("%20", "").lower()
        target_filename = normalize_filename(filename)
        
        # Truy xuất trực tiếp từ bảng documents thay vì FAISS
        res = supabase.table("documents").select("content, metadata").eq("metadata->>user_id", user_id).eq("metadata->>notebook_id", notebook_id).execute()
        
        original_chunks = []
        for doc in res.data:
            doc_filename = normalize_filename(doc.get("metadata", {}).get("filename", ""))
            doc_page = doc.get("metadata", {}).get("page")
            if doc_filename == target_filename and doc_page == page:
                original_chunks.append(doc.get("content", ""))

        if not original_chunks: 
            return {"status": "success", "data": "Không trích xuất được lý thuyết chi tiết cho trang này."}

        raw_theory = "\n\n...\n\n".join(original_chunks)
        format_prompt = f"Đây là văn bản thô. HÃY LÀM 2 VIỆC:\n1. Định dạng lại bằng Markdown (kẻ bảng, in đậm).\n2. TUYỆT ĐỐI GIỮ NGUYÊN 100% THÔNG TIN. KHÔNG THÊM BỚT CHỮ.\n\nVĂN BẢN THÔ:\n{raw_theory}"
        formatted_theory = call_groq(format_prompt, is_chat_mode=False, temp=0.1)
        
        return {"status": "success", "data": formatted_theory}
    except Exception as e: 
        return {"status": "error", "data": f"Lỗi trích xuất: {str(e)}"}